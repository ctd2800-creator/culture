"""Generate LangGraph workflow presentation for KB AI 데이터 리터러시."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "KB_AI_데이터리터러시_LangGraph_워크플로우.pptx"

KB_GOLD = RGBColor(0xFF, 0xB8, 0x00)
KB_DARK = RGBColor(0x1A, 0x1A, 0x2E)
KB_BLUE = RGBColor(0x00, 0x4B, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)

APP_NAME = "KB AI 데이터 리터러시"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, KB_DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = KB_GOLD
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(8.4), Inches(1.6))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(17)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, KB_BLUE)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, sub_bullets=None, font_size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = KB_DARK
        p.space_after = Pt(10)
        p.level = 0

    if sub_bullets:
        for item in sub_bullets:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = GRAY
            p.level = 1
            p.space_after = Pt(6)


def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD

    cols = len(headers)
    nrows = len(rows) + 1
    row_h = min(0.55, 5.5 / nrows)
    table_shape = slide.shapes.add_table(
        nrows, cols, Inches(0.4), Inches(1.1), Inches(9.2), Inches(row_h * nrows)
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = KB_BLUE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = KB_DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def add_flow_slide(prs, title, lines, use_mono=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD

    flow_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(5.8))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xFC)
    flow_box.line.color.rgb = KB_BLUE

    tf = flow_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)

    font_name = "Consolas" if use_mono else "맑은 고딕"
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(13 if not use_mono else 11)
        p.font.color.rgb = KB_DARK
        p.space_after = Pt(3)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        APP_NAME,
        "LangGraph 관점의 워크플로우\n질문 분석 라우터 + 목적별 에이전트 노드 + 공유 State",
    )

    add_content_slide(
        prs,
        "한 줄 요약",
        [
            "Culture 앱은 단일 LangGraph StateGraph 로 동작합니다.",
            "모든 요청은 analyze 노드에서 intent를 판별한 뒤,",
            "해당 전문 에이전트 노드로 분기하고 reply 노드에서 응답을 조립합니다.",
            "",
            "핵심 패턴: START → analyze → (conditional) → 작업 노드 → reply → END",
        ],
        font_size=17,
    )

    add_flow_slide(
        prs,
        "LangGraph 그래프 구조 (노드 · 엣지)",
        [
            "START",
            "  ↓",
            "analyze  (질문 분석 에이전트)",
            "  ↓  route_after_analyze (intent 기반 분기)",
            "  ├─ inst1_extract ──→ fetch_inst1 ──→ format_inst1 ──┐",
            "  ├─ inst1_table_prompt ──→ table_prompt ─────────────┤",
            "  ├─ inst1_aggregate_prompt → aggregate_prompt ───────┤",
            "  ├─ inst1_column_desc ───→ column_desc ──────────────┤",
            "  ├─ inst1_data_summary ──→ data_summary ─────────────┤",
            "  ├─ inst1_chart ─────────→ chart ────────────────────┤",
            "  └─ general_chat ────────→ general ──────────────────┤",
            "                                                       ↓",
            "                                                    reply",
            "                                                       ↓",
            "                                                     END",
        ],
        use_mono=True,
    )

    add_section_slide(prs, "공유 State — CultureState")

    add_table_slide(
        prs,
        "CultureState 주요 필드",
        ["필드", "역할"],
        [
            ("user_message", "현재 사용자 입력"),
            ("question_analysis", "analyze 결과 (intent, tables, month, group_by 등)"),
            ("inst1_data", "조회 결과 rows (테이블별 dict)"),
            ("inst1_queries", "생성된 SQL"),
            ("inst1_column_orders", "UI 컬럼 표시 순서"),
            ("reply / summary", "최종 텍스트 응답"),
            ("chart_specs", "Chart.js 차트 스펙"),
            ("excel_export", "엑셀 다운로드 payload"),
            ("report_export", "PDF 보고서 payload"),
            ("pending_aggregate", "집계 follow-up 세션 상태"),
            ("pending_chart", "차트 follow-up 세션 상태"),
        ],
        col_widths=[3.0, 6.0],
    )

    add_content_slide(
        prs,
        "State 병합 방식",
        [
            "각 노드는 State 전체가 아니라 변경분 dict만 반환합니다.",
            "LangGraph가 반환값을 기존 State에 merge 합니다.",
            "예: fetch_inst1 → { inst1_data, inst1_queries, ... }",
            "예: format_inst1 → { reply, excel_export, pending_chart, ... }",
        ],
        font_size=17,
    )

    add_section_slide(prs, "라우팅 — route_after_analyze")

    add_table_slide(
        prs,
        "intent → 다음 노드",
        ["intent", "다음 노드", "하는 일"],
        [
            ("inst1_table_prompt", "table_prompt", "테이블명만 있을 때 추천 질문 제안"),
            ("inst1_aggregate_prompt", "aggregate_prompt", "집계 전 컬럼 선택 안내"),
            ("inst1_column_desc", "column_desc", "테이블 컬럼 설명"),
            ("inst1_data_summary", "data_summary", "Bedrock 요약 + PDF 보고서"),
            ("inst1_extract", "fetch_inst1 → format_inst1", "SQL 조회 + 결과 포맷 + 엑셀"),
            ("inst1_chart", "chart", "직전 집계 결과로 막대 차트"),
            ("general_chat", "general", "일반 대화 (Bedrock)"),
        ],
        col_widths=[2.4, 2.2, 4.4],
    )

    add_section_slide(prs, "에이전트 노드 상세")

    add_content_slide(
        prs,
        "① analyze — 질문 분석 에이전트",
        [
            "Bedrock + 규칙 기반(analyze_question)으로 intent·테이블·필터 추출",
            "출력: question_analysis (JSON 구조)",
            "",
            "멀티턴 지원:",
            "· Flask 세션의 pending_aggregate, pending_chart를 함께 전달",
            "· 예: 집계 조회 후 「차트 그려줘」 → inst1_chart 로 라우팅",
        ],
        font_size=16,
    )

    add_content_slide(
        prs,
        "② fetch_inst1 → ③ format_inst1 (데이터 조회 파이프라인)",
        [
            "【fetch_inst1】 extract_inst1_data() 실행",
            "· Aurora INST1 테이블 SELECT / aggregate / join_aggregate",
            "· inst1_data, inst1_queries, inst1_column_orders 를 State에 기록",
            "",
            "【format_inst1】 추출 결과 포맷",
            "· format_inst1_reply() 로 텍스트 응답 생성",
            "· 집계 성공 시 pending_chart 저장 → 다음 턴 차트 가능",
            "· 조회 성공 시 excel_export payload 생성",
        ],
        font_size=15,
    )

    add_content_slide(
        prs,
        "④ chart / data_summary / ⑤ reply",
        [
            "【chart】 pending_chart 기반 Chart.js spec 생성 + PDF 보고서",
            "【data_summary】 DB 샘플 조회 후 Bedrock 요약 + PDF 보고서",
            "",
            "【reply】 최종 응답 조립 노드",
            "· 앞 노드의 reply, chart_specs, excel_export 등을 그대로 전달",
            "· 에러 시 [호출 에이전트: ...] 배너와 함께 실패 메시지 반환",
            "· 모든 경로의 공통 종료점 (reply → END)",
        ],
        font_size=16,
    )

    add_section_slide(prs, "멀티턴 시나리오")

    add_flow_slide(
        prs,
        "집계 → 차트 2턴 대화 예시",
        [
            "[1턴] 사용자: 「그룹고객기본정보 성별구분별 집계」",
            "  → analyze → inst1_extract",
            "  → fetch_inst1 → format_inst1",
            "  → pending_chart 저장",
            "  → follow_up: 「조회한 집계 데이터로 차트를 그려드릴까요?」",
            "",
            "[2턴] 사용자: 「차트 그려줘」",
            "  → analyze (pending_chart 참조) → inst1_chart",
            "  → chart → chart_specs + report_export",
            "  → reply → 화면에 그래프 + PDF 보고서 버튼",
        ],
    )

    add_content_slide(
        prs,
        "세션 메모리 (Flask Session)",
        [
            "LangGraph 한 invoke = 한 턴 (그래프 내부 메모리 없음)",
            "턴 간 상태는 Flask session으로 이어갑니다.",
            "",
            "· pending_aggregate — 집계 컬럼 선택 대화",
            "· pending_chart — 집계 후 차트 follow-up",
            "· chat_history — 대화 이력 + inst1_data + 차트/엑셀 메타",
        ],
        font_size=17,
    )

    add_section_slide(prs, "Flask 연동")

    add_table_slide(
        prs,
        "API와 LangGraph 실행",
        ["API", "LangGraph 사용 방식"],
        [
            ("POST /api/chat", "run_workflow() → executor.invoke()"),
            ("POST /api/chat/stream", "executor.stream(updates) — 노드별 진행 표시"),
            ("culture_workflow.py", "build_culture_graph().compile() — 그래프 정의"),
            ("culture_inst1_agents.py", "analyze_question, extract_inst1_data 등 에이전트 로직"),
        ],
        col_widths=[3.2, 5.8],
    )

    add_table_slide(
        prs,
        "LangGraph 개념 ↔ Culture 앱 매핑",
        ["LangGraph 개념", "Culture 앱 구현"],
        [
            ("Graph", "StateGraph(CultureState) 단일 DAG"),
            ("Node", "analyze, fetch_inst1, format_inst1, chart, general 등 10개"),
            ("Edge", "대부분 고정 edge, analyze만 conditional"),
            ("State", "TypedDict 기반 공유 컨텍스트"),
            ("Router", "route_after_analyze (intent 기반)"),
            ("Sub-pipeline", "fetch_inst1 → format_inst1 (조회 전용 2단계)"),
            ("Memory", "LangGraph 내부 X → Flask session"),
            ("LLM", "Amazon Bedrock (analyze, general, data_summary)"),
            ("Action", "SQL 조회, 차트, 엑셀/PDF export (노드 내부 함수)"),
        ],
        col_widths=[2.8, 6.2],
    )

    add_content_slide(
        prs,
        "정리",
        [
            "질문 분석 라우터 + 목적별 전문 에이전트 노드 + 공유 State + reply 종료",
            "",
            "데이터 조회만 2단계 파이프라인(fetch → format)을 거칩니다.",
            "나머지 intent는 단일 노드에서 처리 후 reply로 합류합니다.",
            "",
            "구현 파일: culture_workflow.py (그래프), culture_inst1_agents.py (에이전트)",
        ],
        font_size=17,
    )

    add_title_slide(prs, "감사합니다", "KB AI 데이터 리터러시 · LangGraph Workflow")

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
