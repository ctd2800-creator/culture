"""Generate KB AI 데이터 리터러시 workflow presentation (non-expert friendly)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "KB_FinAgent_Culture_워크플로우_쉬운설명.pptx"

# KB-style colors
KB_GOLD = RGBColor(0xFF, 0xB8, 0x00)
KB_DARK = RGBColor(0x1A, 0x1A, 0x2E)
KB_BLUE = RGBColor(0x00, 0x4B, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, KB_DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = KB_GOLD
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.4))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, KB_BLUE)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, sub_bullets=None, font_size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = KB_DARK
        p.space_after = Pt(10)
        p.level = 0

    if sub_bullets:
        for item in sub_bullets:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = GRAY
            p.level = 1
            p.space_after = Pt(6)


def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD

    cols = len(headers)
    nrows = len(rows) + 1
    row_h = min(0.55, 5.5 / nrows)
    table_shape = slide.shapes.add_table(
        nrows, cols, Inches(0.4), Inches(1.1), Inches(9.2), Inches(row_h * nrows)
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = KB_BLUE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = KB_DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def add_flow_slide(prs, title, lines, use_mono=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD

    flow_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(5.8))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xFC)
    flow_box.line.color.rgb = KB_BLUE

    tf = flow_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)

    font_name = "Consolas" if use_mono else "맑은 고딕"
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font_name
        p.font.size = Pt(14 if not use_mono else 12)
        p.font.color.rgb = KB_DARK
        p.space_after = Pt(4)


APP_NAME = "KB AI 데이터 리터러시"


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 표지 ──
    add_title_slide(
        prs,
        APP_NAME,
        "그룹고객 데이터를 한국어로 질문하고\n답을 받는 AI 분석 채팅 서비스",
    )

    # ── 이 앱이 하는 일 ──
    add_content_slide(
        prs,
        "이 앱이 하는 일",
        [
            "사용자가 평소 말하듯 한국어로 질문합니다.",
            "시스템이 질문의 종류를 먼저 파악합니다.",
            "질문에 맞는 AI 에이전트가 데이터를 조회·정리·설명합니다.",
            "결과는 채팅 화면에 표, 그래프, 저장 버튼으로 보여줍니다.",
        ],
        font_size=18,
    )

    # ── 식당 비유 ──
    add_content_slide(
        prs,
        "전체 흐름 — 쉬운 비유",
        [
            "① 손님(사용자)이 주문합니다 → 질문을 입력합니다.",
            "② 안내 데스크가 주문 내용을 확인합니다 → 어떤 종류의 질문인지 판단합니다.",
            "③ 각 주방이 맡은 요리만 만듭니다 → 컬럼 설명, 집계, 차트 등 해당 기능만 실행합니다.",
            "④ 웨이터가 한꺼번에 서빙합니다 → 답변·표·그래프를 정리해 화면에 표시합니다.",
        ],
        font_size=17,
    )

    add_flow_slide(
        prs,
        "한눈에 보는 처리 순서",
        [
            "사용자가 질문 입력",
            "        ↓",
            "  「이 질문은 어떤 종류인가?」 판단",
            "        ↓",
            "  ┌──────────────────────────────────────┐",
            "  │  테이블 이름만?  → 추천 질문 보여주기              │",
            "  │  컬럼 설명?      → 각 항목 의미 설명                │",
            "  │  데이터 요약?    → AI 요약 + PDF 보고서             │",
            "  │  데이터 조회?    → 표 + 엑셀 저장                  │",
            "  │  집계(합계 등)?  → 표 + 엑셀 + 차트 추천           │",
            "  │  차트?           → 그래프 + PDF 보고서             │",
            "  │  일반 대화?      → 일반 답변                       │",
            "  └──────────────────────────────────────┘",
            "        ↓",
            "  답변을 정리해 화면에 표시",
        ],
    )

    add_section_slide(prs, "분석할 수 있는 데이터")

    add_table_slide(
        prs,
        "분석 가능한 테이블 (3종)",
        ["데이터 이름", "어떤 내용인가"],
        [
            ("그룹고객기본정보", "그룹 고객의 기본 신상·속성 정보"),
            ("그룹고객거래기본", "그룹 고객의 거래·금액 관련 정보"),
            ("그룹고객분석인스턴스목록", "분석에 사용되는 인스턴스 목록"),
        ],
        col_widths=[3.2, 5.8],
    )

    add_section_slide(prs, "질문 종류별로 이렇게 답합니다")

    add_table_slide(
        prs,
        "질문 예시와 시스템 동작",
        ["사용자가 이렇게 물으면", "시스템이 이렇게 합니다"],
        [
            ("「그룹고객거래기본」 (이름만)", "이 테이블로 할 수 있는 질문 4가지를 추천"),
            ("「컬럼 설명해줘」", "각 컬럼(항목)이 무슨 뜻인지 설명"),
            ("「데이터 요약해줘」", "샘플 조회 후 AI 요약 + PDF 보고서 버튼"),
            ("「거래금액 합계 보여줘」", "집계 조회 → 표 + 엑셀 저장 버튼"),
            ("「차트 그려줘」", "막대그래프 + PDF 보고서 버튼 (엑셀 버튼 없음)"),
            ("일반적인 대화", "데이터 조회 없이 일반 답변"),
        ],
        col_widths=[3.8, 5.2],
    )

    add_content_slide(
        prs,
        "① 테이블 이름만 입력한 경우",
        [
            "예: 「그룹고객거래기본」",
            "",
            "시스템이 「이 테이블로 무엇을 해볼까요?」 하고",
            "아래와 같은 추천 질문을 보여줍니다.",
        ],
        sub_bullets=[
            "· 컬럼(항목) 설명 보기",
            "· 데이터 내용 요약하기",
            "· 금액·건수 등 집계하기",
            "· 결과를 그래프로 보기",
        ],
        font_size=17,
    )

    add_content_slide(
        prs,
        "② 컬럼 설명 / ③ 데이터 요약",
        [
            "【컬럼 설명】 예: 「그룹고객거래기본의 컬럼」",
            "→ DB에 등록된 항목 설명을 읽어, 각 컬럼이 무엇을 의미하는지 알려줍니다.",
            "",
            "【데이터 요약】 예: 「그룹고객거래기본 데이터 요약해줘」",
            "→ 샘플 데이터를 AI가 요약하고, 「보고서」 버튼으로 PDF 저장 가능",
        ],
        font_size=16,
    )

    add_section_slide(prs, "저장 기능 — 엑셀 · PDF 보고서")

    add_table_slide(
        prs,
        "에이전트별 저장 버튼",
        ["호출된 에이전트", "화면에 보이는 버튼", "저장 위치"],
        [
            ("데이터 추출 (조회·집계)", "엑셀 저장", "culture\\exports\\"),
            ("데이터 요약", "보고서 (PDF)", "culture\\reports\\"),
            ("차트", "보고서 (PDF)", "culture\\reports\\"),
        ],
        col_widths=[3.0, 2.5, 3.5],
    )

    add_content_slide(
        prs,
        "저장 기능 — 쉬운 설명",
        [
            "【엑셀 저장】 데이터 추출 에이전트가 조회한 표 데이터를 .xlsx로 저장",
            "  · 일반 조회, 집계, 조인 집계 결과 모두 저장 가능",
            "  · 차트 화면에서는 엑셀 버튼이 나오지 않음",
            "",
            "【보고서 (PDF)】 차트·요약 에이전트 결과를 PDF로 저장",
            "  · 요약 텍스트 + (차트인 경우) 막대그래프 포함",
            "  · 브라우저 다운로드 + culture\\reports 폴더에 동시 저장",
        ],
        font_size=16,
    )

    add_section_slide(prs, "집계와 차트 — 2단계 대화")

    add_content_slide(
        prs,
        "④ 집계(합계·건수 등) — 단계별 흐름",
        [
            "1단계: 사용자가 「합계 보여줘」처럼 집계를 요청합니다.",
            "2단계: 시스템이 「어떤 항목으로 집계할까요?」라고 되묻습니다.",
            "3단계: 사용자가 항목(예: 거래금액)을 알려줍니다.",
            "4단계: 시스템이 데이터를 집계해 표로 보여줍니다.",
            "5단계: 「엑셀 저장」 버튼으로 조회 데이터를 파일로 저장할 수 있습니다.",
            "6단계: 「차트도 그려드릴까요?」 추천 질문이 나타납니다.",
        ],
        font_size=17,
    )

    add_content_slide(
        prs,
        "⑤ 차트(그래프) — 단계별 흐름",
        [
            "1단계: 집계 결과가 나온 뒤, 추천 질문을 누르거나 차트를 요청합니다.",
            "2단계: 시스템이 방금 조회한 집계 결과를 기억해 둡니다.",
            "3단계: 그 데이터를 막대그래프로 그려 화면에 표시합니다.",
            "4단계: 「보고서」 버튼으로 PDF 파일을 생성·저장합니다.",
            "",
            "※ 엑셀 저장은 집계 조회 단계에서, PDF는 차트·요약 단계에서 제공됩니다.",
        ],
        font_size=17,
    )

    add_section_slide(prs, "화면에서 보이는 것")

    add_content_slide(
        prs,
        "답변이 화면에 나오는 순서",
        [
            "1. 어떤 에이전트가 처리했는지 안내",
            "     (예: [호출 에이전트: 데이터 추출 에이전트])",
            "2. 설명 문장",
            "3. 생성된 SQL (있을 경우)",
            "4. 조회 결과 표",
            "5. 엑셀 저장 / 보고서(PDF) 버튼 (해당 에이전트일 때)",
            "6. 추천 질문 (예: 차트 그리기)",
            "7. 막대그래프 (차트 요청 시)",
        ],
        font_size=17,
    )

    add_content_slide(
        prs,
        "데이터를 가져올 때의 기본 규칙",
        [
            "기본적으로 「그룹회사코드」 조건만 적용해 데이터를 가져옵니다.",
            "사용자가 질문에 「기준년월」을 넣었을 때만, 그 조건도 함께 적용합니다.",
            "월별로 묶어 집계할 때는, 불필요하게 데이터 범위를 좁히지 않도록",
            "조회 조건과 집계 기준을 구분해 처리합니다.",
        ],
        font_size=17,
    )

    add_content_slide(
        prs,
        "꼭 기억할 4가지",
        [
            "① 질문 종류를 먼저 파악한 뒤, 맞는 AI 에이전트만 실행합니다.",
            "② 집계 → 차트는 2단계 대화로 연결됩니다.",
            "③ 엑셀 저장 = 데이터 추출, PDF 보고서 = 차트·요약 에이전트",
            "④ 화면에는 「처리 안내 → 결과 → 저장·다음 행동」 순으로 표시됩니다.",
        ],
        font_size=19,
    )

    add_title_slide(prs, "감사합니다", APP_NAME)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
