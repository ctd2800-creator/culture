"""채팅 대화 기반 PPTX 보고서 생성 — KB 스타일(흰 배경·노란 포인트·미니멀)."""

from __future__ import annotations

import re
from datetime import datetime, timezone
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from culture_pdf_s3 import _render_chart_png

KB_YELLOW = RGBColor(0xFF, 0xCC, 0x00)
KB_BROWN_DARKER = RGBColor(0x2C, 0x24, 0x19)
KB_BROWN = RGBColor(0x5C, 0x4B, 0x3C)
KB_TEXT_MUTED = RGBColor(0x4A, 0x40, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

KOREAN_FONT = "Malgun Gothic"
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

MARGIN_X = 0.72
BODY_WIDTH = 8.56
TITLE_TOP = 0.52
BODY_TOP = 1.35
BODY_BOTTOM = 7.08
BODY_FONT_PT = 11
BODY_LINE_PT = 14
LINES_PER_SLIDE = 14
WRAP_CHARS = 40

_AGENT_BANNER_RE = re.compile(r"^\[호출 에이전트:[^\]]+\]\s*$")


def _normalize_report_text(text: str) -> str:
    """보고서 본문 — 채팅 원문을 그대로 유지 (에이전트 배너 포함)."""
    return (text or "").strip()


def _report_export_text(item: dict[str, Any]) -> str:
    report = item.get("report_export") or {}
    if not isinstance(report, dict):
        return ""
    return (report.get("content") or report.get("summary") or "").strip()


def _assistant_message_text(item: dict[str, Any]) -> str:
    """답변 본문 — content·report_export 중 더 완전한 원문 사용."""
    content = (item.get("content") or "").strip()
    export_text = _report_export_text(item)
    if len(export_text) > len(content):
        return _normalize_report_text(export_text)
    if content:
        return _normalize_report_text(content)
    return _normalize_report_text(export_text)


def _is_agent_banner_line(line: str) -> bool:
    return bool(_AGENT_BANNER_RE.match((line or "").strip()))


def ascii_report_filename(prefix: str = "culture_report", *, ext: str = "pptx") -> str:
    safe = re.sub(r"[^A-Za-z0-9._-]", "_", (prefix or "culture_report").strip()) or "culture_report"
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    ext = (ext or "pptx").lstrip(".")
    return f"{safe}_{ts}.{ext}"


# 보고서 대상 분석 에이전트 (데이터 요약·외부 환경 분석)
_REPORT_AGENT_NAMES = frozenset({"inst1_data_summary", "inst1_external_insight"})

# 보고서에서 제거할 메타 머리부 ([분석] 블록 bullet)
_ANALYSIS_META_PREFIXES = (
    "- 대상:",
    "- 집계 기준:",
    "- 데이터 건수:",
    "- 기준년월:",
    "- 그룹회사코드:",
)
_SUMMARY_HEADING = "집계 핵심 요약"
_REPORT_TITLE_RE = re.compile(r"^\s*보고서\s*제목\s*[:：]\s*(.+?)\s*$")

# 분석 소제목 (번호 접두사 포함/미포함 모두 인식)
_SECTION_HEADINGS = ("집계 핵심 요약", "외부 환경 연결", "시사점·제언", "시사점", "제언")
_SECTION_HEADING_RE = re.compile(
    r"^\s*\d+\s*[.)]?\s*(집계\s*핵심\s*요약|외부\s*환경\s*연결|시사점[·\s]*제언)\s*$"
)


def _is_section_heading(line: str) -> bool:
    """'1. 집계 핵심 요약' 같은 분석 소제목 라인 여부."""
    if _SECTION_HEADING_RE.match(line or ""):
        return True
    stripped = (line or "").strip()
    return stripped in _SECTION_HEADINGS


def _extract_report_title(content: str) -> str:
    """분석 본문에서 '보고서 제목: ...' 라인을 찾아 제목 텍스트 반환."""
    for line in (content or "").split("\n"):
        m = _REPORT_TITLE_RE.match(line)
        if m:
            return m.group(1).strip()
    return ""


def _clean_analysis_content(content: str) -> str:
    """분석 본문에서 에이전트 배너·[분석] 메타·안내 문구 제거 → '집계 핵심 요약'부터."""
    lines = (content or "").split("\n")
    for idx, line in enumerate(lines):
        norm = re.sub(r"[\d.)\s]", "", line)
        if norm.startswith("집계핵심요약"):
            lines = lines[idx:]
            break
    else:
        cleaned: list[str] = []
        for line in lines:
            stripped = line.strip()
            if _is_agent_banner_line(stripped):
                continue
            if stripped == "[분석]":
                continue
            if _REPORT_TITLE_RE.match(line):
                continue
            if stripped.startswith(_ANALYSIS_META_PREFIXES):
                continue
            cleaned.append(line)
        lines = cleaned
    out: list[str] = []
    for line in lines:
        if "보고서 버튼" in line and "생성" in line:
            continue
        out.append(line)
    return "\n".join(out).strip()


def _split_analysis_head_tail(content: str) -> tuple[str, str]:
    """본문을 '외부 환경 연결' 기준으로 head(요약)·tail(외부 환경~)로 분리."""
    lines = (content or "").split("\n")
    for idx, line in enumerate(lines):
        norm = re.sub(r"[#*\-0-9.()\[\]\s]", "", line)
        if norm.startswith("외부환경연결") and len(norm) <= len("외부환경연결") + 2:
            head = "\n".join(lines[:idx]).rstrip()
            tail = "\n".join(lines[idx:]).strip()
            return head, tail
    return content, ""


def chat_history_has_data(history: list[dict[str, Any]]) -> bool:
    for item in history:
        if item.get("role") != "assistant":
            continue
        if _assistant_message_text(item):
            return True
        if item.get("charts"):
            return True
        inst1 = item.get("inst1_data") or {}
        if any(rows for rows in inst1.values() if rows):
            return True
    return False


def _merge_assistant_item(
    server_item: dict[str, Any] | None,
    client_item: dict[str, Any] | None,
) -> dict[str, Any]:
    server_item = dict(server_item or {})
    client_item = dict(client_item or {})
    merged = dict(server_item) if len(server_item) >= len(client_item) else dict(client_item)
    if not merged:
        return {"role": "assistant", "content": ""}

    server_text = _assistant_message_text(server_item) if server_item else ""
    client_text = _assistant_message_text(client_item) if client_item else ""
    merged["content"] = server_text if len(server_text) >= len(client_text) else client_text

    server_report = server_item.get("report_export") if server_item else None
    client_report = client_item.get("report_export") if client_item else None
    if isinstance(server_report, dict) or isinstance(client_report, dict):
        pick = server_report if len(_report_export_text(server_item or {})) >= len(
            _report_export_text(client_item or {})
        ) else client_report
        if isinstance(pick, dict) and pick:
            merged["report_export"] = dict(pick)

    for key in ("charts", "inst1_data", "inst1_column_orders", "inst1_result_labels", "inst1_queries"):
        server_val = server_item.get(key) if server_item else None
        client_val = client_item.get(key) if client_item else None
        if server_val:
            merged[key] = server_val
        elif client_val:
            merged[key] = client_val

    merged["role"] = "assistant"
    return merged


def merge_chat_histories(
    server_history: list[dict[str, Any]] | None,
    client_history: list[dict[str, Any]] | None,
) -> list[dict[str, Any]]:
    """서버 세션·클라이언트 기록 병합 — 긴 답변(분석 에이전트) 누락 방지."""
    server_history = list(server_history or [])
    client_history = list(client_history or [])
    if not server_history:
        return client_history
    if not client_history:
        return server_history

    merged: list[dict[str, Any]] = []
    count = max(len(server_history), len(client_history))
    for idx in range(count):
        server_item = server_history[idx] if idx < len(server_history) else {}
        client_item = client_history[idx] if idx < len(client_history) else {}
        role = server_item.get("role") or client_item.get("role") or "user"
        if role == "assistant":
            merged.append(_merge_assistant_item(server_item, client_item))
            continue
        pick = server_item if server_item else client_item
        merged.append(
            {
                "role": "user",
                "content": (pick.get("content") or "").strip(),
            }
        )
    return merged


def apply_report_patch(
    history: list[dict[str, Any]],
    report_patch: dict[str, Any] | None,
) -> list[dict[str, Any]]:
    """보고서 버튼이 가리키는 분석 답변 본문을 history에 반영."""
    if not history or not isinstance(report_patch, dict):
        return history

    patch_text = (report_patch.get("content") or report_patch.get("summary") or "").strip()
    if not patch_text:
        return history

    patch_agent = (report_patch.get("agent") or "").strip()
    out = [dict(item) for item in history]

    for idx in range(len(out) - 1, -1, -1):
        item = out[idx]
        if item.get("role") != "assistant":
            continue
        item_agent = ""
        report = item.get("report_export") or {}
        if isinstance(report, dict):
            item_agent = (report.get("agent") or "").strip()
        if patch_agent and patch_agent not in {item_agent, ""}:
            if patch_agent not in (_assistant_message_text(item) or ""):
                continue
        existing = _assistant_message_text(item)
        if len(patch_text) >= len(existing):
            item["content"] = patch_text
            if isinstance(report, dict):
                merged_report = dict(report)
            else:
                merged_report = {}
            merged_report.setdefault("agent", patch_agent or item_agent)
            merged_report["content"] = patch_text
            item["report_export"] = merged_report
            out[idx] = item
        break

    return out


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _apply_font(run, *, size: int = BODY_FONT_PT, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = KOREAN_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    return slide


def _add_top_accent_line(slide) -> None:
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = KB_YELLOW
    line.line.fill.background()


def _add_title_underline(slide, y: float, width: float = 1.6, x: float | None = None) -> None:
    left = Inches(x if x is not None else MARGIN_X)
    shape = slide.shapes.add_shape(1, left, Inches(y), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = KB_YELLOW
    shape.line.fill.background()



def _hard_wrap_line(line: str, width: int = WRAP_CHARS) -> list[str]:
    text = line.rstrip()
    if not text:
        return [""]
    parts: list[str] = []
    while text:
        if len(text) <= width:
            parts.append(text)
            break
        chunk = text[:width]
        break_at = max(chunk.rfind(" "), chunk.rfind("，"), chunk.rfind(","), chunk.rfind("。"))
        if break_at > width // 4:
            parts.append(text[:break_at].rstrip())
            text = text[break_at:].lstrip()
        else:
            parts.append(text[:width])
            text = text[width:].lstrip()
    return parts


def _explode_to_display_lines(body: str) -> list[str]:
    """본문을 슬라이드에 실제로 그릴 줄 단위로 펼침 (누락·잘림 방지)."""
    display: list[str] = []
    for line in (body or "").split("\n"):
        if not line.strip():
            display.append("")
            continue
        display.extend(_hard_wrap_line(line))
    return display


def _paginate_lines(lines: list[str], *, per_page: int = LINES_PER_SLIDE) -> list[list[str]]:
    if not lines:
        return [[]]
    return [lines[i : i + per_page] for i in range(0, len(lines), per_page)]


def _fill_body_text(tf, lines: list[str]) -> None:
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_bottom = tf.margin_top = Pt(2)
    tf.margin_left = tf.margin_right = Pt(2)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line:
            p.text = line
        else:
            p.text = ""
        p.line_spacing = BODY_LINE_PT
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        if not p.runs:
            p.add_run()
        if _is_agent_banner_line(line):
            _apply_font(p.runs[0], size=BODY_FONT_PT + 1, bold=True, color=KB_BROWN_DARKER)
        elif line.lstrip().startswith(("-", "•", "·")):
            _apply_font(p.runs[0], color=KB_BROWN)
        else:
            _apply_font(p.runs[0], color=KB_BROWN_DARKER)


def _add_slide_title(slide, title: str, *, size: int = 19) -> None:
    box = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(TITLE_TOP), Inches(BODY_WIDTH), Inches(0.68)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    _apply_font(p.runs[0], size=size, bold=True, color=KB_BROWN_DARKER)
    _add_title_underline(slide, TITLE_TOP + 0.74, width=1.25)


def _add_body_slides(prs: Presentation, title: str, body: str) -> None:
    display_lines = _explode_to_display_lines(body)
    pages = _paginate_lines(display_lines)
    for page_idx, page_lines in enumerate(pages):
        slide = _new_slide(prs)
        _add_top_accent_line(slide)
        slide_title = title if page_idx == 0 else f"{title} (계속 {page_idx + 1}/{len(pages)})"
        _add_slide_title(slide, slide_title)
        body_box = slide.shapes.add_textbox(
            Inches(MARGIN_X),
            Inches(BODY_TOP),
            Inches(BODY_WIDTH),
            Inches(BODY_BOTTOM - BODY_TOP),
        )
        _fill_body_text(body_box.text_frame, page_lines)


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = _new_slide(prs)
    _add_top_accent_line(slide)

    brand = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(1.05), Inches(BODY_WIDTH), Inches(0.4))
    bp = brand.text_frame.paragraphs[0]
    bp.text = "KB 금융그룹"
    bp.alignment = PP_ALIGN.CENTER
    _apply_font(bp.runs[0], size=12, bold=True, color=KB_TEXT_MUTED)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(8.2), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    _apply_font(p.runs[0], size=32, bold=True, color=KB_BROWN_DARKER)

    _add_title_underline(slide, 3.35, width=2.0, x=4.0)

    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.9), Inches(3.65), Inches(8.2), Inches(1.6))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(subtitle.split("\n")):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = line
            p2.alignment = PP_ALIGN.CENTER
            p2.space_after = Pt(4)
            _apply_font(p2.runs[0], size=15, color=KB_TEXT_MUTED)


def _add_section_slide(prs: Presentation, index: int, title: str) -> None:
    """질문 제목 — 길면 여러 슬라이드로 분할."""
    display = _explode_to_display_lines(title)
    pages = _paginate_lines(display, per_page=10)
    for page_idx, page_lines in enumerate(pages):
        slide = _new_slide(prs)
        _add_top_accent_line(slide)
        label = f"Q{index}" if page_idx == 0 else f"Q{index} (계속)"
        _add_slide_title(slide, label, size=22)

        body_box = slide.shapes.add_textbox(
            Inches(MARGIN_X),
            Inches(2.05),
            Inches(BODY_WIDTH),
            Inches(BODY_BOTTOM - 2.05),
        )
        tf = body_box.text_frame
        tf.word_wrap = False
        for i, line in enumerate(page_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.line_spacing = Pt(20)
            p.space_after = Pt(4)
            if p.runs:
                _apply_font(p.runs[0], size=20, bold=True, color=KB_BROWN_DARKER)


def _add_chart_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    title = str(spec.get("title") or "차트")
    slide = _new_slide(prs)
    _add_top_accent_line(slide)
    _add_slide_title(slide, title, size=18)
    chart_top = BODY_TOP + 0.05
    try:
        png = _render_chart_png(spec)
        slide.shapes.add_picture(
            BytesIO(png),
            Inches(MARGIN_X),
            Inches(chart_top),
            width=Inches(BODY_WIDTH),
        )
    except Exception:
        body = slide.shapes.add_textbox(
            Inches(MARGIN_X), Inches(chart_top), Inches(BODY_WIDTH), Inches(0.8)
        )
        p = body.text_frame.paragraphs[0]
        p.text = f"(차트 렌더 실패: {title})"
        _apply_font(p.runs[0], color=KB_TEXT_MUTED)


def _order_columns(rows: list[dict[str, Any]], preferred: list[str] | None) -> list[str]:
    if preferred:
        cols = [c for c in preferred if rows and c in rows[0]]
        if cols:
            return cols
    if not rows:
        return []
    cols = list(rows[0].keys())
    if "고객수" in cols:
        return [c for c in cols if c != "고객수"] + ["고객수"]
    return cols


def _format_inst1_table_lines(
    *,
    label: str,
    rows: list[dict[str, Any]],
    column_order: list[str] | None,
) -> list[str]:
    if not rows:
        return []
    cols = _order_columns(rows, column_order)
    lines = [f"{label}  ({len(rows)}건)", ""]
    header = "  ·  ".join(cols)
    lines.append(header)
    lines.append("─" * min(len(header), 80))
    for row in rows:
        cells = [str(row.get(c, "") if row.get(c) is not None else "") for c in cols]
        row_text = "  ·  ".join(cells)
        lines.extend(_hard_wrap_line(row_text, width=WRAP_CHARS))
    return lines


def _assistant_sections(history: list[dict[str, Any]]) -> list[tuple[str, dict[str, Any]]]:
    sections: list[tuple[str, dict[str, Any]]] = []
    pending_user = ""
    answer_idx = 0
    for item in history:
        role = item.get("role")
        if role == "user":
            pending_user = (item.get("content") or "").strip()
            continue
        if role != "assistant":
            continue
        answer_idx += 1
        title = pending_user or f"답변 {answer_idx}"
        sections.append((title, item))
        pending_user = ""
    return sections


def build_chat_report_ppt_bytes(history: list[dict[str, Any]]) -> bytes:
    if not chat_history_has_data(history):
        raise ValueError("보고서로 저장할 내용이 없습니다.")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    generated = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    _add_title_slide(
        prs,
        "KB AI 데이터 리터러시",
        f"분석 보고서\n생성 시각: {generated}",
    )

    sections = _assistant_sections(history)
    if not sections:
        raise ValueError("보고서로 저장할 답변이 없습니다.")

    summary_lines = [f"{i + 1}. {title}" for i, (title, _) in enumerate(sections)]
    _add_body_slides(prs, "대화 요약", "\n".join(summary_lines))

    for idx, (question, item) in enumerate(sections, start=1):
        _add_section_slide(prs, idx, question)

        content = _assistant_message_text(item)
        if content:
            slide_title = "분석 결과"
            for line in content.split("\n"):
                if _is_agent_banner_line(line):
                    slide_title = line.strip().strip("[]")
                    break
            _add_body_slides(prs, slide_title, content)

        for spec in item.get("charts") or []:
            if isinstance(spec, dict):
                _add_chart_slide(prs, spec)

        inst1_data = item.get("inst1_data") or {}
        labels = item.get("inst1_result_labels") or {}
        column_orders = item.get("inst1_column_orders") or {}
        for table_name, rows in inst1_data.items():
            if not rows:
                continue
            label = str(labels.get(table_name) or table_name)
            table_lines = _format_inst1_table_lines(
                label=label,
                rows=rows,
                column_order=column_orders.get(table_name),
            )
            _add_body_slides(prs, "조회 데이터", "\n".join(table_lines))

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def latest_analysis_section(
    history: list[dict[str, Any]]
) -> tuple[str, dict[str, Any]] | None:
    """가장 최근 호출된 분석 에이전트 답변(질문 제목, 항목) 반환."""
    sections = _assistant_sections(history)
    if not sections:
        return None
    for title, item in reversed(sections):
        report = item.get("report_export") or {}
        agent = (report.get("agent") or "").strip() if isinstance(report, dict) else ""
        if agent in _REPORT_AGENT_NAMES:
            return title, item
    return sections[-1]


def _doc_set_korean_font(run, *, size: int = 11, bold: bool = False,
                         color: tuple[int, int, int] | None = None) -> None:
    from docx.oxml.ns import qn
    from docx.shared import Pt as DocPt, RGBColor as DocRGB

    run.font.name = KOREAN_FONT
    run.font.size = DocPt(size)
    run.font.bold = bold
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = rpr.makeelement(qn("w:rFonts"), {})
        rpr.append(rfonts)
    rfonts.set(qn("w:eastAsia"), KOREAN_FONT)
    if color is not None:
        run.font.color.rgb = DocRGB(*color)


def _doc_add_paragraph(doc, text: str, *, size: int = 11, bold: bool = False,
                       color: tuple[int, int, int] | None = None,
                       space_after: int = 4):
    from docx.shared import Pt as DocPt

    p = doc.add_paragraph()
    p.paragraph_format.space_after = DocPt(space_after)
    p.paragraph_format.space_before = DocPt(0)
    run = p.add_run(text or "")
    _doc_set_korean_font(run, size=size, bold=bold, color=color)
    return p


def _doc_add_body_text(doc, body: str) -> None:
    """분석 본문을 줄 단위로 추가 (소제목·배너·불릿 강조)."""
    for line in (body or "").split("\n"):
        stripped = line.rstrip()
        if not stripped.strip():
            doc.add_paragraph()
            continue
        if _is_section_heading(stripped):
            _doc_add_paragraph(
                doc, stripped.strip(), size=15, bold=True,
                color=(0x2C, 0x24, 0x19), space_after=6,
            )
        elif _is_agent_banner_line(stripped):
            _doc_add_paragraph(
                doc, stripped, size=13, bold=True,
                color=(0x2C, 0x24, 0x19), space_after=6,
            )
        elif stripped.lstrip().startswith(("-", "•", "·")):
            _doc_add_paragraph(doc, stripped, size=11, color=(0x5C, 0x4B, 0x3C))
        else:
            _doc_add_paragraph(doc, stripped, size=11, color=(0x2C, 0x24, 0x19))


def _doc_add_inst1_table(doc, *, rows: list[dict[str, Any]],
                         column_order: list[str] | None) -> None:
    if not rows:
        return
    cols = _order_columns(rows, column_order)
    if not cols:
        return
    table = doc.add_table(rows=1, cols=len(cols))
    try:
        table.style = "Light Grid Accent 1"
    except Exception:
        pass
    hdr = table.rows[0].cells
    for i, col in enumerate(cols):
        hdr[i].text = ""
        run = hdr[i].paragraphs[0].add_run(str(col))
        _doc_set_korean_font(run, size=10, bold=True, color=(0x2C, 0x24, 0x19))
    for row in rows:
        cells = table.add_row().cells
        for i, col in enumerate(cols):
            val = row.get(col)
            cells[i].text = ""
            run = cells[i].paragraphs[0].add_run(
                "" if val is None else str(val)
            )
            _doc_set_korean_font(run, size=10, color=(0x4A, 0x40, 0x38))
    doc.add_paragraph()


def _doc_add_chart(doc, spec: dict[str, Any]) -> None:
    from docx.shared import Inches as DocInches

    try:
        # 원형·도넛은 범례가 조각 라벨을 보여주므로 집계 항목명을 캡션으로 표시
        chart_type = str(spec.get("type") or "").strip().lower()
        if chart_type in ("pie", "doughnut"):
            datasets = spec.get("datasets") or []
            measure_name = str(datasets[0].get("label") or "").strip() if datasets else ""
            if measure_name:
                _doc_add_paragraph(
                    doc, measure_name, size=11, bold=True,
                    color=(0x2C, 0x24, 0x19), space_after=2,
                )
        png = _render_chart_png(spec, show_title=False)
        doc.add_picture(BytesIO(png), width=DocInches(6.2))
    except Exception:
        pass


def build_chat_report_docx_bytes(history: list[dict[str, Any]]) -> bytes:
    """가장 최근 분석 에이전트 결과만으로 Word(.docx) 보고서 생성.

    구성: 제목 → 집계 핵심 요약(본문) → 표 → 차트 → 외부 환경 연결~ (꼬리).
    분석 질문·에이전트 배너·[분석] 메타·생성 시각·구역 머리글은 표시하지 않는다.
    """
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt as DocPt

    section = latest_analysis_section(history)
    if section is None:
        raise ValueError("보고서로 저장할 분석 결과가 없습니다.")
    _question, item = section

    raw_content = _assistant_message_text(item)
    report_title = _extract_report_title(raw_content)
    content = _clean_analysis_content(raw_content)
    head, tail = _split_analysis_head_tail(content)
    charts = [s for s in (item.get("charts") or []) if isinstance(s, dict)]
    inst1_data = item.get("inst1_data") or {}
    has_table = any(rows for rows in inst1_data.values() if rows)
    if not content and not charts and not has_table:
        raise ValueError("보고서로 저장할 분석 결과가 없습니다.")

    doc = Document()
    for sec in doc.sections:
        sec.top_margin = sec.bottom_margin = DocPt(54)
        sec.left_margin = sec.right_margin = DocPt(54)

    if report_title:
        title_p = doc.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _doc_set_korean_font(title_p.add_run(report_title),
                             size=20, bold=True, color=(0x2C, 0x24, 0x19))
        doc.add_paragraph()

    if head:
        _doc_add_body_text(doc, head)

    column_orders = item.get("inst1_column_orders") or {}
    if has_table:
        doc.add_paragraph()
        for table_name, rows in inst1_data.items():
            if not rows:
                continue
            _doc_add_inst1_table(
                doc,
                rows=rows,
                column_order=column_orders.get(table_name),
            )

    if charts:
        doc.add_paragraph()
        for spec in charts:
            _doc_add_chart(doc, spec)

    if tail:
        doc.add_paragraph()
        _doc_add_body_text(doc, tail)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()
