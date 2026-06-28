"""Generate one-slide KB FinAgent Culture tech stack — 6 shape boxes."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "KB_FinAgent_Culture_주요기술_한장_도형.pptx"

KB_GOLD = RGBColor(0xFF, 0xB8, 0x00)
KB_DARK = RGBColor(0x1A, 0x1A, 0x2E)
KB_BLUE = RGBColor(0x00, 0x4B, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

BOXES = [
    {
        "label": "프론트엔드",
        "tech": "HTML·CSS  |  culture_chat.js\nChart.js  |  SSE",
        "role": "채팅 UI · 표 · 그래프 · 스트리밍",
    },
    {
        "label": "백엔드",
        "tech": "Python 3.12  |  Flask\nWerkzeug",
        "role": "로그인 · 세션 · REST/SSE API",
    },
    {
        "label": "AI",
        "tech": "LangGraph  |  AWS Bedrock\nboto3  ·  Claude Sonnet 4.5",
        "role": "질문 분기 · SQL 생성 · 요약 · 답변",
    },
    {
        "label": "데이터베이스",
        "tech": "PostgreSQL (Aurora)\npsycopg2  ·  INST1 × 3",
        "role": "그룹고객 데이터 조회 · 집계",
    },
    {
        "label": "배포",
        "tech": "Vercel  |  AWS IAM",
        "role": "서버리스 배포 · Bedrock 인증",
    },
    {
        "label": "선택 기능",
        "tech": "S3  |  fpdf2  |  matplotlib",
        "role": "요약 PDF 생성 · 저장",
    },
]


def add_category_box(slide, left, top, width, height, label, tech, role):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = KB_BLUE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p_label = tf.paragraphs[0]
    p_label.text = label
    p_label.font.size = Pt(15)
    p_label.font.bold = True
    p_label.font.color.rgb = KB_BLUE
    p_label.alignment = PP_ALIGN.CENTER
    p_label.space_after = Pt(6)

    p_tech = tf.add_paragraph()
    p_tech.text = tech
    p_tech.font.size = Pt(11)
    p_tech.font.bold = True
    p_tech.font.color.rgb = KB_DARK
    p_tech.alignment = PP_ALIGN.CENTER
    p_tech.space_after = Pt(8)

    p_role = tf.add_paragraph()
    p_role.text = f"→ {role}"
    p_role.font.size = Pt(10)
    p_role.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p_role.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = KB_DARK
    bar.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.65))
    tp = title.text_frame.paragraphs[0]
    tp.text = "KB FinAgent (Culture) — 주요 기술"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = KB_GOLD
    tp.alignment = PP_ALIGN.CENTER

    # 3 × 2 grid of boxes
    margin_x = 0.35
    margin_y = 1.15
    gap_x = 0.22
    gap_y = 0.22
    cols, rows = 3, 2
    usable_w = 10 - margin_x * 2
    usable_h = 7.5 - margin_y - 0.55  # leave room for bottom flow line
    box_w = (usable_w - gap_x * (cols - 1)) / cols
    box_h = (usable_h - gap_y * (rows - 1)) / rows

    for i, box in enumerate(BOXES):
        col = i % cols
        row = i // cols
        left = Inches(margin_x + col * (box_w + gap_x))
        top = Inches(margin_y + row * (box_h + gap_y))
        add_category_box(
            slide,
            left,
            top,
            Inches(box_w),
            Inches(box_h),
            box["label"],
            box["tech"],
            box["role"],
        )

    # bottom flow line
    flow = slide.shapes.add_textbox(Inches(0.35), Inches(6.95), Inches(9.3), Inches(0.45))
    fp = flow.text_frame.paragraphs[0]
    fp.text = (
        "흐름:  브라우저 → Flask → LangGraph → [PostgreSQL | Bedrock] → 응답  "
        "|  핵심: culture_app · culture_workflow · culture_inst1_agents · culture_chat.js"
    )
    fp.font.size = Pt(9.5)
    fp.font.color.rgb = KB_DARK
    fp.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
