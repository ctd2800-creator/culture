import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_kb_presentation():
    prs = Presentation()
    # 16:9 와이드스크린 비율 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 테마 색상 설정 (KB 컬러 분위기 반영)
    COLOR_PRIMARY = RGBColor(74, 74, 74)   # 챠콜 (본문/타이틀)
    COLOR_ACCENT = RGBColor(255, 184, 28)  # KB 골드 (강조색)
    COLOR_MUTED = RGBColor(128, 128, 128)  # 그레이 (부제/설명)

    # 헬퍼 함수: 텍스트 박스 생성 및 스타일 적용
    def add_title(slide, text):
        tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = 'Malgun Gothic'
        p.font.color.rgb = COLOR_PRIMARY
        return tx_box

    # -------------------------------------------------------------
    # Slide 1: 타이틀
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # 빈 슬라이드
    slide1 = prs.slides.add_slide(slide_layout)
    
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.833), Inches(2.5))
    tf = title_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "KB AI 데이터 리터러시"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    p1.font.name = 'Malgun Gothic'
    
    p2 = tf.add_paragraph()
    p2.text = "자연어 기반 데이터 추출, 외부 정보 융합 및 보고서 제작 원스톱 시스템"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_MUTED
    p2.font.name = 'Malgun Gothic'
    p2.space_before = Pt(20)

    # -------------------------------------------------------------
    # Slide 2: 배경 및 Pain Point
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    add_title(slide2, "데이터 추출 대기 및 보고서 수제작의 지체")
    
    content_box2 = slide2.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.833), Inches(4.5))
    tf2 = content_box2.text_frame
    points2 = [
        "• 데이터 확보 대기: 수작업 쿼리 작성 등으로 인해 데이터 추출 요청 후 전달까지 평균 1~2일 소요",
        "• 보고서 제작 비효율: 로우 데이터의 엑셀 정렬, 수기 차트 제작, 문장 작성 및 레이아웃 구성 등 수제작 리소스 과다",
        "• 통합 분석의 한계: 내부 데이터와 외부 시장 동향을 결합한 입체적 분석의 어려움"
    ]
    for i, pt in enumerate(points2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.name = 'Malgun Gothic'
        p.space_after = Pt(24)

    # -------------------------------------------------------------
    # Slide 3: 핵심 아이디어 및 독창성
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    add_title(slide3, "데이터 분석부터 보고서 제작까지의 지능형 자동화 에이전트")
    
    content_box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.833), Inches(4.5))
    tf3 = content_box3.text_frame
    points3 = [
        "• 자율형 데이터 분석: 자연어 기반으로 최적의 쿼리를 도출하고 내부 금융 데이터를 추출 및 분석",
        "• 외부 정보 융합: 경제/정책/시장 동향 정보를 실시간으로 연동하여 데이터 분석의 깊이와 통찰력을 보강",
        "• 목적별 산출물 제작: 엑셀 원본 데이터, 실시간 동적 차트, 분석 해설이 포함된 완성형 보고서 직접 제작"
    ]
    for i, pt in enumerate(points3):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.name = 'Malgun Gothic'
        p.space_after = Pt(24)

    # -------------------------------------------------------------
    # Slide 4: 기술적 타당성 (아키텍처)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    add_title(slide4, "분석 신뢰성을 보장하는 기술 아키텍처")
    
    content_box4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
    tf4 = content_box4.text_frame
    tf4.word_wrap = True
    
    # 한 줄 요약
    p_sum = tf4.paragraphs[0]
    p_sum.text = "요약: KB AI 데이터 리터러시 앱은 OpenSearch 데이터 사전으로 질문 의도를 정확히 파악하고, LangGraph로 추출·집계·차트·분석을 단계별 전문 에이전트에 위임하며, AWS Bedrock(AI 추론) + Aurora(정형 데이터)를 결합해 데이터 추출부터 해석·보고서 제작까지 자동화한 지능형 파이프라인입니다."
    p_sum.font.size = Pt(14)
    p_sum.font.italic = True
    p_sum.font.name = 'Malgun Gothic'
    p_sum.font.color.rgb = COLOR_MUTED
    p_sum.space_after = Pt(20)
    
    # 3단계 기술 파이프라인
    pipeline_text = (
        "1. 의도 파악 — 데이터 사전 기반 검색 (OpenSearch)\n"
        "   - 의미 좌표 변환(임베딩), 데이터 사전 매핑(k-NN), 지능형 힌트 주입을 통해 데이터 정확도 확보\n\n"
        "2. 워크플로우 제어 — 모듈형 에이전트 협업 (LangGraph)\n"
        "   - 11개 노드가 질문 분석부터 응답 조립까지 순차 제어, Intent 기반 전문 에이전트로 업무 라우팅\n\n"
        "3. 실행·해석 — 하이브리드 엔진 (Bedrock + Aurora)\n"
        "   - 규칙 기반 SQL 생성 및 Aurora 조회, Bedrock 기반 분석·인사이트 도출 및 보고서 자동 생성"
    )
    p_pipe = tf4.add_paragraph()
    p_pipe.text = pipeline_text
    p_pipe.font.size = Pt(15)
    p_pipe.font.name = 'Malgun Gothic'
    p_pipe.space_after = Pt(20)

    # 차별성
    p_rob = tf4.add_paragraph()
    p_rob.text = "※ 기술적 차별성: 규칙 기반 SQL + 스키마 힌트 (안정성) | 모듈형 아키텍처 (확장성) | 세션 스냅샷 (멀티턴 연속 분석)"
    p_rob.font.size = Pt(14)
    p_rob.font.bold = True
    p_rob.font.name = 'Malgun Gothic'
    p_rob.font.color.rgb = COLOR_PRIMARY

    # -------------------------------------------------------------
    # Slide 5: 실시간 라이브 시연
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    add_title(slide5, "[Live Demo] 자연어 한 번으로 끝내는 데이터 리포트")
    
    content_box5 = slide5.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(4.5))
    tf5 = content_box5.text_frame
    steps = [
        "Step 1. [질문 입력] \"지난달 우리 그룹의 마케팅 성성과를 이번 달 금리 정책 변화와 연계해서 분석해줘.\"",
        "Step 2. [실시간 처리] 질문 의도 파악 → 데이터 사전 매핑 → 데이터 추출 → 시각화 프로세스 시연",
        "Step 3. [결과물 확인] Raw Data 엑셀 파일 다운로드, 동적 차트 출력, AI 서술형 인사이트가 담긴 완성형 PPT 확인"
    ]
    for i, step in enumerate(steps):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.name = 'Malgun Gothic'
        p.space_after = Pt(28)

    # -------------------------------------------------------------
    # Slide 6: 결론 및 향후 계획
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    add_title(slide6, "검증 완료 및 향후 데이터 자율화 방향")
    
    content_box6 = slide6.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(4.5))
    tf6 = content_box6.text_frame
    points6 = [
        "• 검증 성과: 기존 1~2일 소요되던 데이터 수집/가공/보고 작업을 단 3분으로 단축, 실무형 보고서 자동 제작 실증 완료",
        "• 향후 계획: 사내 AI 포털 연동을 통한 전사 분석 환경 조성 및 자연어 기반 보고 자동화 아키텍처 확대 적용"
    ]
    for i, pt in enumerate(points6):
        p = tf6.paragraphs[0] if i == 0 else tf6.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.name = 'Malgun Gothic'
        p.space_after = Pt(28)

    # 파일 저장
    output_fn = "KB_AI_Data_Literacy_Presentation.pptx"
    prs.save(output_fn)
    print(f"[성공] 슬라이드가 {output_fn} 파일로 생성되었습니다.")

if __name__ == "__main__":
    create_kb_presentation()
