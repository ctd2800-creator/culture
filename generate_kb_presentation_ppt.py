"""KB AI 데이터 리터러시 발표용 PPT — VBA CreateKBPPT()와 동일 구성."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "KB_AI_데이터리터러시_발표.pptx"

KB_YELLOW = RGBColor(0xFF, 0xCC, 0x00)
KB_BROWN_DARKER = RGBColor(0x2C, 0x24, 0x19)
KB_BROWN = RGBColor(0x5C, 0x4B, 0x3C)
KB_TEXT_MUTED = RGBColor(0x4A, 0x40, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

KOREAN_FONT = "Malgun Gothic"
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
MARGIN_X = 0.72
BODY_WIDTH = 8.56


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _font(run, *, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = KOREAN_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = KB_YELLOW
    line.line.fill.background()
    return slide


def _underline(slide, y: float, *, width: float = 1.6, x: float | None = None) -> None:
    shape = slide.shapes.add_shape(
        1,
        Inches(x if x is not None else MARGIN_X),
        Inches(y),
        Inches(width),
        Inches(0.04),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = KB_YELLOW
    shape.line.fill.background()


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = _blank_slide(prs)

    brand = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(1.05), Inches(BODY_WIDTH), Inches(0.4))
    bp = brand.text_frame.paragraphs[0]
    bp.text = "KB 금융그룹"
    bp.alignment = PP_ALIGN.CENTER
    _font(bp.runs[0], size=12, bold=True, color=KB_TEXT_MUTED)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(8.2), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    _font(p.runs[0], size=34, bold=True, color=KB_BROWN_DARKER)
    _underline(slide, 3.2, width=2.2, x=3.9)

    box2 = slide.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(8.2), Inches(2.4))
    tf = box2.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(subtitle.split("\n")):
        p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p2.text = line or " "
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(6)
        if not p2.runs:
            p2.add_run()
        if not line.strip():
            continue
        _font(
            p2.runs[0],
            size=16 if i == 0 else 14,
            color=KB_TEXT_MUTED if i else KB_BROWN,
        )


def add_content_slide(prs: Presentation, title: str, body: str) -> None:
    slide = _blank_slide(prs)

    title_box = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(0.52), Inches(BODY_WIDTH), Inches(0.85)
    )
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.word_wrap = True
    _font(tp.runs[0], size=22, bold=True, color=KB_BROWN_DARKER)
    _underline(slide, 1.28, width=1.4)

    body_box = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(1.55), Inches(BODY_WIDTH), Inches(5.5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.line_spacing = Pt(18)
        p.space_after = Pt(4)
        if not p.runs:
            p.add_run()
        stripped = line.strip()
        if stripped and stripped[0].isdigit() and "." in stripped[:3]:
            _font(p.runs[0], size=15, bold=True, color=KB_BROWN_DARKER)
        elif stripped.startswith("-"):
            _font(p.runs[0], size=13, color=KB_BROWN)
        elif stripped.startswith("Step"):
            _font(p.runs[0], size=14, bold=True, color=KB_BROWN_DARKER)
        elif stripped.startswith("["):
            _font(p.runs[0], size=14, bold=True, color=KB_BROWN_DARKER)
        else:
            _font(p.runs[0], size=13, color=KB_BROWN_DARKER)


def main() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(
        prs,
        "KB AI 데이터 리터러시",
        "자연어 기반 데이터 추출, 외부 정보 융합 및 보고서 제작 원스톱 시스템\n\n"
        "발표자: 컬쳐리더팀 최승윤 차장, 김홍철 차장, 김태훈 차장",
    )

    add_content_slide(
        prs,
        "데이터 추출 대기 및 보고서 수제작의 지체",
        "1. 데이터 확보 대기\n"
        " - 수작업 쿼리 작성으로 평균 1~2일 소요\n\n"
        "2. 보고서 제작 비효율\n"
        " - 엑셀 정렬, 차트 제작, PPT 레이아웃 구성 등 수작업 과다\n\n"
        "3. 통합 분석의 한계\n"
        " - 내부 데이터와 외부 시장(경제/정책) 동향 결합의 어려움",
    )

    add_content_slide(
        prs,
        "데이터 분석부터 보고서 제작까지 지능형 자동화",
        "1. 자율형 데이터 분석\n"
        " - 자연어 기반 최적 쿼리 도출 및 데이터 추출\n\n"
        "2. 외부 정보 융합\n"
        " - 실시간 경제/정책 연동으로 인사이트 보강\n\n"
        "3. 목적별 산출물 제작\n"
        " - Excel 다운로드, 동적 Chart, PPT 보고서 자동 레이아웃 구성",
    )

    add_content_slide(
        prs,
        "지능형 분석 및 보고서 제작 파이프라인",
        "Step 1. 데이터 구조 사전 학습 (메타데이터 추출)\n"
        "Step 2. 지능형 지도 구축 (벡터 임베딩/의미 지도)\n"
        "Step 3. 의도 기반 검색 (OpenSearch/k-NN 매칭)\n"
        "Step 4. 복합 분석 제어 (LangGraph 자가 보정 루프)",
    )

    add_content_slide(
        prs,
        "프로세스 단축에 따른 정량적 업무 효율화",
        "[AS-IS vs TO-BE 비교]\n\n"
        "- 데이터 추출: 1~2일 대기 -> 1분 (즉시 실행)\n"
        "- 엑셀 가공: 1시간 수작업 -> 즉시 다운로드\n"
        "- 시각화 분석: 1시간 소요 -> 실시간 차트 노출\n"
        "- PPT 제작: 2시간 소요 -> 1분 (생산성 90% 극대화)",
    )

    add_content_slide(
        prs,
        "검증 완료 및 향후 데이터 자율화 방향",
        "1. 검증 성과\n"
        " - 외부 정보 융합 및 실무형 PPT 직접 제작 실증 완료\n"
        " - 1~2일 소요 업무를 단 3분으로 단축\n\n"
        "2. 향후 계획\n"
        " - 사내 AI 포털 연동 추진\n"
        " - 전사 데이터 분석 및 보고 자동화 통합 업무 환경 완성",
    )

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    main()
